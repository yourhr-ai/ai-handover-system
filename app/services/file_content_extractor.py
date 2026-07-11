from pathlib import Path

MAX_SUMMARY_LENGTH = 800
TRUNCATION_SUFFIX = "...(이하 생략)"
MAX_EXCEL_DATA_ROWS = 9
MAX_PDF_FIRST_PAGE_CHARS = 500


def extract_file_summary(absolute_path: str) -> str | None:
    extension = Path(absolute_path).suffix.lower()

    try:
        if extension in (".xlsx", ".xls"):
            summary = _extract_excel_summary(absolute_path)
        elif extension == ".docx":
            summary = _extract_docx_summary(absolute_path)
        elif extension == ".pdf":
            summary = _extract_pdf_summary(absolute_path)
        elif extension == ".pptx":
            summary = _extract_pptx_summary(absolute_path)
        else:
            return None
    except Exception:
        return None

    if not summary:
        return None

    return _truncate_summary(summary)


def _truncate_summary(text: str) -> str:
    text = text.strip()
    if len(text) <= MAX_SUMMARY_LENGTH:
        return text
    return text[:MAX_SUMMARY_LENGTH] + TRUNCATION_SUFFIX


def _extract_excel_summary(absolute_path: str) -> str | None:
    import openpyxl

    workbook = openpyxl.load_workbook(absolute_path, read_only=True, data_only=True)
    try:
        sheet = workbook.worksheets[0]
        rows = list(
            sheet.iter_rows(
                min_row=1,
                max_row=1 + MAX_EXCEL_DATA_ROWS,
                values_only=True,
            )
        )
    finally:
        workbook.close()

    if not rows:
        return None

    header_text = _format_excel_row(rows[0])
    data_lines = [
        line
        for line in (_format_excel_row(row) for row in rows[1:])
        if line
    ]

    parts = []
    if header_text:
        parts.append(f"헤더: {header_text}")
    if data_lines:
        parts.append("데이터: " + " / ".join(data_lines))

    return "\n".join(parts) if parts else None


def _format_excel_row(row: tuple) -> str:
    return ", ".join(str(value) for value in row if value is not None)


def _extract_docx_summary(absolute_path: str) -> str | None:
    from docx import Document

    document = Document(absolute_path)
    paragraphs = [
        paragraph for paragraph in document.paragraphs if paragraph.text.strip()
    ]
    if not paragraphs:
        return None

    first_paragraphs = [paragraph.text.strip() for paragraph in paragraphs[:2]]
    last_paragraph = (
        [paragraphs[-1].text.strip()] if len(paragraphs) > len(first_paragraphs) else []
    )
    heading_paragraphs = [
        paragraph.text.strip()
        for paragraph in paragraphs
        if paragraph.style is not None
        and paragraph.style.name.startswith("Heading")
    ]

    parts = []
    if first_paragraphs:
        parts.append("시작 문단: " + " / ".join(first_paragraphs))
    if last_paragraph:
        parts.append("마지막 문단: " + last_paragraph[0])
    if heading_paragraphs:
        parts.append("제목: " + " / ".join(heading_paragraphs))

    return "\n".join(parts) if parts else None


def _extract_pdf_summary(absolute_path: str) -> str | None:
    import pdfplumber

    with pdfplumber.open(absolute_path) as pdf:
        if not pdf.pages:
            return None
        text = (pdf.pages[0].extract_text() or "").strip()

    if not text:
        return None

    return text[:MAX_PDF_FIRST_PAGE_CHARS]


def _extract_pptx_summary(absolute_path: str) -> str | None:
    from pptx import Presentation

    presentation = Presentation(absolute_path)
    titles = [
        slide.shapes.title.text.strip()
        for slide in presentation.slides
        if slide.shapes.title is not None and slide.shapes.title.text.strip()
    ]

    if not titles:
        return None

    return "슬라이드 제목: " + " / ".join(titles)
