import hashlib
import base64
from collections import deque
from concurrent.futures import (
    FIRST_COMPLETED,
    ProcessPoolExecutor,
    ThreadPoolExecutor,
    TimeoutError,
    as_completed,
    wait,
)
import json
import logging
import math
import os
import re
import shutil
import tempfile
import time
import warnings
import zipfile
from datetime import date, datetime
from html import unescape
from pathlib import Path
from typing import Any, Callable

from app.services.analysis_result import AnalysisResult, AnalyzedFile
from app.services.parallel_file_runner import (
    recommended_file_workers,
    run_process_items_with_timeout,
    terminate_process_executor,
)


EMBEDDING_MODEL = "text-embedding-3-small"
EMBEDDING_BATCH_SIZE = 200
EMBEDDING_BATCH_MAX_CHARS = 200_000
FILE_EXTRACTION_MAX_WORKERS = recommended_file_workers()
DEFAULT_FILE_EXTRACTION_TIMEOUT_SECONDS = 10
SLOW_FILE_EXTRACTION_TIMEOUT_SECONDS = 15
COST_ESTIMATION_TIMEOUT_SECONDS = 5
EMBEDDING_MAX_WORKERS = 2
MAX_EXCEL_EMBED_ROWS = 100
MAX_EXCEL_EMBED_COLUMNS = 20


class RagPackageCancelled(Exception):
    pass


def _raise_if_cancelled(cancel_check: Callable[[], bool] | None) -> None:
    if cancel_check is not None and cancel_check():
        raise RagPackageCancelled()
EMBEDDING_RETRY_ATTEMPTS = 5
EMBEDDING_REQUEST_DELAY_SECONDS = 0.2
PDF_EXTRACT_TIMEOUT_SECONDS = 12
RAG_TEXT_EXTRACTION_EXTENSIONS = {
    ".docx", ".pptx", ".ppt", ".txt", ".md", ".xlsx", ".xls", ".hwp", ".hwpx"
}
SUPPORTED_EMAIL_ATTACHMENT_EXTENSIONS = RAG_TEXT_EXTRACTION_EXTENSIONS
SYSTEM_FILE_NAMES = {"thumbs.db", "desktop.ini", ".ds_store"}
COMMON_REMOVAL_KEYWORDS = (
    "confidential",
    "all rights reserved",
    "copyright",
    "©",
    "무단전재",
    "저작권",
    "수신거부",
    "수신 거부",
    "unsubscribe",
    "발송되었습니다",
    "클릭하세요",
    "클릭해주세요",
    "if you don't want to receive",
)
EMAIL_SIGNATURE_KEYWORDS = (
    "회사",
    "직책",
    "부서",
    "전화",
    "연락처",
    "mobile",
    "tel",
    "fax",
    "email",
    "e-mail",
)
_EXTRACTED_TEXT_CACHE: dict[str, str] = {}
_ATTACHMENT_TEXT_CACHE: dict[str, str] = {}
_FILE_HASH_CACHE: dict[tuple[str, int, int], str] = {}
_PDF_TEXT_TIMEOUT_PATHS: set[str] = set()


def filter_files_by_selected_extensions(
    files: list,
    selected_extensions: set[str],
    max_file_size_bytes: int | None = None,
    extension_size_limits: dict[str, int | None] | None = None,
) -> tuple[list, list]:
    """Split folder files by whether their contents should be embedded.

    Every file remains part of the package. Files outside the selected RAG
    whitelist, or larger than ``max_file_size_bytes``, are represented by the
    existing filename-only placeholder. ``None`` preserves the unlimited
    behavior used by older callers.
    """
    normalized = {extension.lower() for extension in selected_extensions}
    normalized_limits = (
        {extension.lower(): limit for extension, limit in extension_size_limits.items()}
        if extension_size_limits is not None
        else None
    )
    content_files = []
    filename_only_files = []
    for item in files:
        analyzed_file = item[0] if isinstance(item, tuple) else item
        extension = Path(analyzed_file.file_name).suffix.lower()
        if normalized_limits is None:
            is_selected = extension in normalized
            is_within_size_limit = (
                max_file_size_bytes is None
                or analyzed_file.size_bytes <= max_file_size_bytes
            )
        else:
            extension_limit = normalized_limits.get(extension, 0)
            is_selected = extension in normalized_limits and extension_limit != 0
            is_within_size_limit = (
                extension_limit is None
                or (isinstance(extension_limit, int) and analyzed_file.size_bytes <= extension_limit)
            )
        target = (
            content_files
            if (
                extension in RAG_TEXT_EXTRACTION_EXTENSIONS
                and is_selected
                and is_within_size_limit
            )
            else filename_only_files
        )
        target.append(item)
    return content_files, filename_only_files


def get_rag_package_candidate_files(
    analysis_result: AnalysisResult,
    folder_paths: list[str] | None = None,
    cancel_check: Callable[[], bool] | None = None,
) -> list[AnalyzedFile]:
    """Return the exact de-duplicated folder-file set used by cost/build.

    This hashes the full content of every candidate file for dedup, so it can
    take a long time on large folders — callers running it off the UI thread
    should pass ``cancel_check`` to allow interrupting between files.
    """
    selected_files, _exclusions = _collect_rag_package_candidate_files(
        analysis_result,
        folder_paths or _split_root_folder_paths(analysis_result),
        cancel_check,
    )
    return [file for file, _absolute_path in selected_files]


def extract_text_for_rag(file_path: str) -> str:
    extension = Path(file_path).suffix.lower()
    if extension in {".txt", ".md", ".csv", ".json", ".xml", ".yaml", ".yml"}:
        return _extract_plain_text_file(file_path)
    if extension == ".docx":
        return _extract_docx_text(file_path)
    if extension == ".pdf":
        return _extract_pdf_text(file_path)
    if extension == ".xlsx":
        return _extract_xlsx_text(file_path)
    if extension == ".xls":
        return _extract_xls_text(file_path)
    if extension == ".pptx":
        return _extract_pptx_text(file_path)
    if extension in {".hwp", ".hwpx"}:
        from src.core.hwp_extractor import HwpExtractor

        try:
            return HwpExtractor().extract(Path(file_path))
        except Exception:
            return ""
    return ""


def _build_unsupported_placeholder_text(
    file_name: str,
    source_path: str,
    modified_at: str,
) -> str:
    return (
        f"[파일명: {file_name}] 이 파일은 자동으로 내용을 읽을 수 없는 형식"
        f"(예: hwp)입니다. 제목을 참고하시고, 세부 내용은 원본 파일을 직접 "
        f"확인하세요. 원본 경로: {source_path}, 수정일시: {modified_at}"
    )


def _build_unsupported_placeholder_chunk_record(
    source_file: str,
    file_name: str,
    source_path: str,
    modified_at: str,
) -> dict:
    return {
        "chunk_text": _build_unsupported_placeholder_text(
            file_name,
            source_path,
            modified_at,
        ),
        "source_file": source_file,
        "chunk_index": 0,
        "metadata": {
            "source_path": source_path,
            "file_name": file_name,
            "modified_at": modified_at,
            "extraction_status": "unsupported_format",
        },
    }


def _build_not_in_whitelist_placeholder_text(
    file_name: str,
    source_path: str,
    modified_at: str,
) -> str:
    return (
        f"[파일명: {file_name}] 이 파일은 지원 대상이 아니라 내용을 자동으로 "
        f"읽지 않았습니다. 제목을 참고하시고, 세부 내용은 원본 파일을 직접 "
        f"확인하세요. 원본 경로: {source_path}, 수정일시: {modified_at}"
    )


def _build_not_in_whitelist_placeholder_chunk_record(
    source_file: str,
    file_name: str,
    source_path: str,
    modified_at: str,
) -> dict:
    return {
        "chunk_text": _build_not_in_whitelist_placeholder_text(
            file_name,
            source_path,
            modified_at,
        ),
        "source_file": source_file,
        "chunk_index": 0,
        "metadata": {
            "source_path": source_path,
            "file_name": file_name,
            "modified_at": modified_at,
            "extraction_status": "not_in_whitelist",
        },
    }


def _extract_plain_text_file(file_path: str) -> str:
    for encoding in ("utf-8", "cp949", "euc-kr"):
        try:
            return Path(file_path).read_text(encoding=encoding, errors="ignore")
        except OSError:
            return ""
        except UnicodeError:
            continue
    return ""


def _extract_docx_text(file_path: str) -> str:
    try:
        from docx import Document

        document = Document(file_path)
        parts = [paragraph.text.strip() for paragraph in document.paragraphs]
        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(part for part in parts if part)
    except Exception:
        return ""


def _extract_pdf_text(file_path: str) -> str:
    executor = ThreadPoolExecutor(max_workers=1)
    future = executor.submit(_extract_pdf_text_without_timeout, file_path)
    try:
        return future.result(timeout=PDF_EXTRACT_TIMEOUT_SECONDS)
    except TimeoutError:
        _PDF_TEXT_TIMEOUT_PATHS.add(str(Path(file_path)))
        future.cancel()
        return ""
    finally:
        executor.shutdown(wait=False, cancel_futures=True)


def _extract_pdf_text_without_timeout(file_path: str) -> str:
    try:
        logging.getLogger("pdfminer").setLevel(logging.ERROR)
        logging.getLogger("pdfplumber").setLevel(logging.ERROR)
        with warnings.catch_warnings():
            warnings.filterwarnings("ignore", message=".*FontBBox.*")
            import pdfplumber

            parts = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text() or ""
                    if text.strip():
                        parts.append(text.strip())
            return "\n\n".join(parts)
    except Exception:
        return ""


def _extract_xlsx_text(file_path: str) -> str:
    try:
        from openpyxl import load_workbook

        value_workbook = load_workbook(file_path, read_only=False, data_only=True)
        formula_workbook = load_workbook(file_path, read_only=False, data_only=False)
        parts: list[str] = []
        try:
            for sheet_index, value_sheet in enumerate(value_workbook.worksheets):
                formula_sheet = formula_workbook.worksheets[sheet_index]
                merged_values = _build_xlsx_merged_value_map(
                    value_sheet,
                    formula_sheet,
                    max_row=min(max(value_sheet.max_row, formula_sheet.max_row), MAX_EXCEL_EMBED_ROWS + 1),
                    max_column=min(max(value_sheet.max_column, formula_sheet.max_column), MAX_EXCEL_EMBED_COLUMNS),
                )

                def get_value(row: int, column: int) -> str:
                    merged = merged_values.get((row, column))
                    if merged is not None:
                        value, formula_value, number_format = merged
                    else:
                        value_cell = value_sheet.cell(row=row, column=column)
                        formula_cell = formula_sheet.cell(row=row, column=column)
                        value = value_cell.value
                        formula_value = formula_cell.value
                        number_format = formula_cell.number_format or value_cell.number_format
                    return _format_excel_value(value, number_format, formula_value)

                parts.extend(
                    _build_excel_sheet_lines(
                        value_sheet.title,
                        max(value_sheet.max_row, formula_sheet.max_row),
                        max(value_sheet.max_column, formula_sheet.max_column),
                        get_value,
                    )
                )
        finally:
            value_workbook.close()
            formula_workbook.close()
        return "\n".join(parts)
    except Exception:
        return ""


def _extract_xls_text(file_path: str) -> str:
    try:
        import xlrd

        workbook = xlrd.open_workbook(file_path, formatting_info=True)
        parts: list[str] = []
        for worksheet in workbook.sheets():
            merged_values: dict[tuple[int, int], tuple[object, str]] = {}
            for row_low, row_high, column_low, column_high in worksheet.merged_cells:
                if row_low >= MAX_EXCEL_EMBED_ROWS + 1 or column_low >= MAX_EXCEL_EMBED_COLUMNS:
                    continue
                anchor = worksheet.cell(row_low, column_low)
                number_format = _get_xls_number_format(workbook, anchor.xf_index)
                for row in range(row_low, min(row_high, MAX_EXCEL_EMBED_ROWS + 1)):
                    for column in range(column_low, min(column_high, MAX_EXCEL_EMBED_COLUMNS)):
                        merged_values[(row, column)] = (anchor.value, number_format)

            def get_value(row: int, column: int) -> str:
                zero_based = (row - 1, column - 1)
                merged = merged_values.get(zero_based)
                if merged is not None:
                    value, number_format = merged
                else:
                    cell = worksheet.cell(*zero_based)
                    value = cell.value
                    number_format = _get_xls_number_format(workbook, cell.xf_index)
                    if cell.ctype == xlrd.XL_CELL_DATE:
                        value = xlrd.xldate_as_datetime(value, workbook.datemode)
                return _format_excel_value(value, number_format)

            parts.extend(
                _build_excel_sheet_lines(
                    worksheet.name,
                    worksheet.nrows,
                    worksheet.ncols,
                    get_value,
                )
            )
        return "\n".join(parts)
    except Exception:
        return ""


def _build_xlsx_merged_value_map(
    value_sheet,
    formula_sheet,
    *,
    max_row: int | None = None,
    max_column: int | None = None,
) -> dict:
    merged_values: dict[tuple[int, int], tuple[object, object, str]] = {}
    for merged_range in formula_sheet.merged_cells.ranges:
        if max_row is not None and merged_range.min_row > max_row:
            continue
        if max_column is not None and merged_range.min_col > max_column:
            continue
        anchor_formula_cell = formula_sheet.cell(
            row=merged_range.min_row,
            column=merged_range.min_col,
        )
        anchor_value_cell = value_sheet.cell(
            row=merged_range.min_row,
            column=merged_range.min_col,
        )
        merged_value = (
            anchor_value_cell.value,
            anchor_formula_cell.value,
            anchor_formula_cell.number_format or anchor_value_cell.number_format,
        )
        for row in range(merged_range.min_row, min(merged_range.max_row, max_row or merged_range.max_row) + 1):
            for column in range(merged_range.min_col, min(merged_range.max_col, max_column or merged_range.max_col) + 1):
                merged_values[(row, column)] = merged_value
    return merged_values


def _build_excel_sheet_lines(
    sheet_name: str,
    max_row: int,
    max_column: int,
    get_value: Callable[[int, int], str],
) -> list[str]:
    lines = [f"[시트: {sheet_name}]"]
    if max_row <= 0 or max_column <= 0:
        return lines

    limited_max_column = min(max_column, MAX_EXCEL_EMBED_COLUMNS)
    headers = [get_value(1, column) for column in range(1, limited_max_column + 1)]
    header_last_column = _last_non_empty_column(headers)
    if header_last_column:
        lines.append("[헤더] " + " | ".join(headers[:header_last_column]))

    last_data_row = min(max_row, MAX_EXCEL_EMBED_ROWS + 1)
    for row in range(2, last_data_row + 1):
        values = [get_value(row, column) for column in range(1, limited_max_column + 1)]
        value_last_column = _last_non_empty_column(values)
        if value_last_column == 0:
            continue
        last_column = max(header_last_column, value_last_column)
        fields = []
        for column in range(1, last_column + 1):
            header = headers[column - 1].strip() or f"열 {column}"
            fields.append(f"{header}: {values[column - 1]}")
        lines.append(f"[행 {row}] " + " | ".join(fields))
    if max_row > last_data_row:
        lines.append(f"[EXCEL_TRUNCATED: total_rows={max_row - 1}, embedded_rows=100]")
        lines.append(f"(이하 생략, 전체 {max_row - 1}행 중 100행만 표시됨)")
    if max_column > limited_max_column:
        lines.append(f"[EXCEL_COLUMNS_TRUNCATED: total_columns={max_column}, embedded_columns=20]")
        lines.append(f"(이하 열 생략, 전체 {max_column}열 중 20열만 표시됨)")
    return lines


def _last_non_empty_column(values: list[str]) -> int:
    for index in range(len(values), 0, -1):
        if values[index - 1] != "":
            return index
    return 0


def _format_excel_value(
    value: object,
    number_format: str = "General",
    formula_value: object | None = None,
) -> str:
    if value is None:
        if isinstance(formula_value, str) and formula_value.startswith("="):
            return "수식 결과를 확인할 수 없습니다"
        return ""
    if isinstance(value, bool):
        return "TRUE" if value else "FALSE"
    if isinstance(value, datetime):
        return value.isoformat(sep=" ", timespec="seconds")
    if isinstance(value, date):
        return value.isoformat()
    if isinstance(value, (int, float)):
        return _format_excel_number(value, number_format)
    return str(value).strip()


def _format_excel_number(value: int | float, number_format: str) -> str:
    format_section = (number_format or "General").split(";", 1)[0]
    if "%" in format_section:
        decimals = _excel_decimal_places(format_section, "%")
        return f"{value * 100:,.{decimals}f}%"

    currency = "₩" if "₩" in format_section or "\\₩" in format_section else ""
    if not currency and "$" in format_section:
        currency = "$"
    grouping = "," in format_section
    decimals = _excel_decimal_places(format_section)
    if currency or grouping:
        return f"{currency}{value:,.{decimals}f}"
    if decimals > 0:
        return f"{value:.{decimals}f}"
    return str(value)


def _excel_decimal_places(number_format: str, stop_character: str = "") -> int:
    section = number_format.split(stop_character, 1)[0] if stop_character else number_format
    match = re.search(r"[0#]\.([0#]+)", section)
    return len(match.group(1)) if match else 0


def _get_xls_number_format(workbook, xf_index: int) -> str:
    try:
        format_key = workbook.xf_list[xf_index].format_key
        return str(workbook.format_map[format_key].format_str)
    except (AttributeError, IndexError, KeyError):
        return "General"


def _extract_pptx_text(file_path: str) -> str:
    try:
        from pptx import Presentation

        presentation = Presentation(file_path)
        parts = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_parts = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text.strip():
                    slide_parts.append(text.strip())
            if slide_parts:
                parts.append(f"[슬라이드 {slide_index}]\n" + "\n".join(slide_parts))
        return "\n\n".join(parts)
    except Exception:
        return ""


def preprocess_text_for_rag(text: str, source_type: str) -> str:
    if source_type == "kakao":
        return text
    if source_type == "email":
        if _looks_like_html(text):
            text = _strip_html_to_text(text)
        text = _remove_lines_with_keywords(
            text,
            (*COMMON_REMOVAL_KEYWORDS, *EMAIL_SIGNATURE_KEYWORDS),
        )
        return _remove_email_signature(_remove_email_reply_chain(text)).strip()
    if source_type == "document":
        return _remove_lines_with_keywords(text, COMMON_REMOVAL_KEYWORDS).strip()
    return text


def split_into_chunks(
    text: str,
    chunk_size: int = 700,
    overlap: int = 100,
) -> list[str]:
    cleaned = text.strip()
    if not cleaned:
        return []

    chunks: list[str] = []
    current = ""
    paragraphs = [p.strip() for p in re.split(r"\n\s*\n+", cleaned) if p.strip()]
    for paragraph in paragraphs:
        paragraph_parts = _split_long_text(paragraph, chunk_size, overlap)
        for part in paragraph_parts:
            if not current:
                current = part
                continue
            candidate = f"{current}\n\n{part}"
            if len(candidate) <= chunk_size:
                current = candidate
            else:
                chunks.append(current)
                prefix = _tail_overlap(current, overlap)
                current = f"{prefix}\n\n{part}" if prefix else part

    if current:
        chunks.append(current)
    return [chunk for chunk in chunks if chunk.strip()]


def split_spreadsheet_into_chunks(
    text: str,
    chunk_size: int = 700,
) -> list[str]:
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if not lines:
        return []

    chunks: list[str] = []
    current_sheet = ""
    current_lines: list[str] = []

    def flush() -> None:
        nonlocal current_lines
        if current_lines and any(not line.startswith("[시트: ") for line in current_lines):
            chunks.append("\n".join(current_lines))
        current_lines = [current_sheet] if current_sheet else []

    for line in lines:
        if line.startswith("[시트: "):
            flush()
            current_sheet = line
            current_lines = [line]
            continue

        candidate_lines = [*current_lines, line]
        if len("\n".join(candidate_lines)) <= chunk_size:
            current_lines = candidate_lines
            continue

        flush()
        candidate_lines = [*current_lines, line]
        if len("\n".join(candidate_lines)) <= chunk_size:
            current_lines = candidate_lines
            continue

        available_size = max(100, chunk_size - len(current_sheet) - 1)
        for part in _split_long_text(line, available_size, overlap=0):
            chunks.append("\n".join(filter(None, (current_sheet, part))))
        current_lines = [current_sheet] if current_sheet else []

    flush()
    return chunks


def _split_text_for_file(text: str, file_name: str) -> list[str]:
    if Path(file_name).suffix.lower() in {".xlsx", ".xls"}:
        return split_spreadsheet_into_chunks(text)
    return split_into_chunks(text)


def embed_chunks(
    chunks: list[str],
    api_key: str,
    progress_callback: Callable[[str, int, int], None] | None = None,
    cancel_check: Callable[[], bool] | None = None,
) -> list[dict]:
    embedded, _failures, _embedding_tokens = _embed_chunks_with_failures(
        chunks,
        api_key,
        progress_callback,
        cancel_check,
    )
    return embedded


def _embed_chunks_with_failures(
    chunks: list[str],
    api_key: str,
    progress_callback: Callable[[str, int, int], None] | None = None,
    cancel_check: Callable[[], bool] | None = None,
) -> tuple[list[dict], list[dict], int]:
    from openai import OpenAI

    client = OpenAI(api_key=api_key)
    batches = _build_embedding_batches(chunks)
    total_batches = len(batches)
    if total_batches == 0:
        return [], [], 0
    _raise_if_cancelled(cancel_check)

    def embed_batch(start: int, batch: list[str]) -> tuple[list[dict], int]:
        _raise_if_cancelled(cancel_check)
        if EMBEDDING_REQUEST_DELAY_SECONDS > 0:
            time.sleep(EMBEDDING_REQUEST_DELAY_SECONDS)
        _raise_if_cancelled(cancel_check)
        response = _create_embedding_with_retry(client, batch)
        _raise_if_cancelled(cancel_check)
        vectors = [item.embedding for item in response.data]
        batch_embedded: list[dict] = []
        for offset, (chunk_text, vector) in enumerate(zip(batch, vectors)):
            batch_embedded.append(
                {
                    "chunk_text": chunk_text,
                    "embedding_vector": vector,
                    "source_file": "",
                    "chunk_index": start + offset,
                }
            )
        usage = getattr(response, "usage", None)
        tokens = int(getattr(usage, "total_tokens", 0) or getattr(usage, "prompt_tokens", 0) or 0)
        return batch_embedded, tokens

    results: list[list[dict] | None] = [None] * total_batches
    failures: list[dict] = []
    embedding_tokens = 0
    completed_batches = 0
    max_workers = min(EMBEDDING_MAX_WORKERS, total_batches)
    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        futures = {
            executor.submit(embed_batch, start, batch): batch_index
            for batch_index, (start, batch) in enumerate(batches)
        }
        for future in as_completed(futures):
            if cancel_check is not None and cancel_check():
                for pending_future in futures:
                    pending_future.cancel()
                executor.shutdown(wait=False, cancel_futures=True)
                raise RagPackageCancelled()
            batch_index = futures[future]
            start, batch = batches[batch_index]
            try:
                batch_result, batch_tokens = future.result()
                results[batch_index] = batch_result
                embedding_tokens += batch_tokens
            except Exception as exc:
                error_message = f"{type(exc).__name__}: {exc}"
                for offset, chunk_text in enumerate(batch):
                    failures.append(
                        {
                            "chunk_index": start + offset,
                            "batch_index": batch_index,
                            "chunk_text_preview": chunk_text[:500],
                            "error": error_message,
                        }
                    )
            completed_batches += 1
            if progress_callback is not None:
                progress_callback("embedding", completed_batches, total_batches)
            _raise_if_cancelled(cancel_check)

    embedded: list[dict] = []
    for batch_result in results:
        if batch_result:
            embedded.extend(batch_result)
    return embedded, failures, embedding_tokens


def _build_embedding_batches(chunks: list[str]) -> list[tuple[int, list[str]]]:
    batches: list[tuple[int, list[str]]] = []
    current_start = 0
    current_batch: list[str] = []
    current_chars = 0
    for index, chunk in enumerate(chunks):
        chunk_chars = len(chunk)
        should_flush = (
            current_batch
            and (
                len(current_batch) >= EMBEDDING_BATCH_SIZE
                or current_chars + chunk_chars > EMBEDDING_BATCH_MAX_CHARS
            )
        )
        if should_flush:
            batches.append((current_start, current_batch))
            current_start = index
            current_batch = []
            current_chars = 0
        current_batch.append(chunk)
        current_chars += chunk_chars
    if current_batch:
        batches.append((current_start, current_batch))
    return batches


def save_rag_package(
    chunks_with_embeddings: list[dict],
    analysis_result: AnalysisResult,
    output_path: str,
    api_key: str,
    source_map: dict[str, dict[str, Any]] | None = None,
    excluded_files: list[dict] | None = None,
    embedding_failures: list[dict] | None = None,
    cancel_check: Callable[[], bool] | None = None,
) -> str:
    zip_path = _get_rag_package_zip_path(Path(output_path))
    zip_path.parent.mkdir(parents=True, exist_ok=True)
    source_map = source_map or {}
    excluded_files = excluded_files or []
    embedding_failures = embedding_failures or []

    with tempfile.TemporaryDirectory(prefix="handover_rag_package_") as temp_dir:
        _raise_if_cancelled(cancel_check)
        package_path = Path(temp_dir) / zip_path.stem
        package_path.mkdir(parents=True, exist_ok=True)

        manifest = {
            "package_name": package_path.name,
            "created_at": datetime.now().isoformat(timespec="seconds"),
            "analysis_id": _build_analysis_id(analysis_result),
            "document_count": len(
                {
                    item.get("source_file", "")
                    for item in chunks_with_embeddings
                    if item.get("source_file", "")
                }
            ),
            "chunk_count": len(chunks_with_embeddings),
            "embedding_failed_chunk_count": len(embedding_failures),
            "embedding_model": EMBEDDING_MODEL,
        }
        (package_path / "manifest.json").write_text(
            json.dumps(manifest, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        with (package_path / "chunks.jsonl").open("w", encoding="utf-8") as file:
            for index, item in enumerate(chunks_with_embeddings):
                _raise_if_cancelled(cancel_check)
                chunk_id = item.get("chunk_id") or f"chunk-{index:06d}"
                file.write(
                    json.dumps(
                        {
                            "chunk_id": chunk_id,
                            "source_file": item.get("source_file", ""),
                            "chunk_text": item.get("chunk_text", ""),
                            "embedding": item.get("embedding_vector", []),
                        },
                        ensure_ascii=False,
                    )
                    + "\n"
                )
        (package_path / "source_map.json").write_text(
            json.dumps(source_map, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        (package_path / "excluded_files.json").write_text(
            json.dumps(excluded_files, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        (package_path / "embedding_failures.json").write_text(
            json.dumps(embedding_failures, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        (package_path / "api_key.dat").write_bytes(base64.b64encode(api_key.encode("utf-8")))
        _raise_if_cancelled(cancel_check)
        _zip_package_folder(package_path, zip_path, cancel_check)
        _raise_if_cancelled(cancel_check)

    return str(zip_path)


def estimate_rag_package_cost(
    analysis_result: AnalysisResult,
    folder_paths: list[str] | None = None,
    parsed_emails: list[dict] | None = None,
    kakao_file_paths: list[str] | None = None,
    cancel_check: Callable[[], bool] | None = None,
    progress_callback: Callable[[str, int, int], None] | None = None,
    selected_extensions: set[str] | None = None,
    max_file_size_bytes: int | None = None,
    extension_size_limits: dict[str, int | None] | None = None,
    embedding_unit_cost_per_1k: float | None = None,
) -> dict[str, int | float]:
    _raise_if_cancelled(cancel_check)
    selected_files, _file_exclusions = _collect_rag_package_candidate_files(
        analysis_result,
        folder_paths or _split_root_folder_paths(analysis_result),
        cancel_check,
    )
    content_files, filename_only_files = filter_files_by_selected_extensions(
        selected_files,
        selected_extensions
        if selected_extensions is not None
        else set(RAG_TEXT_EXTRACTION_EXTENSIONS),
        max_file_size_bytes,
        extension_size_limits,
    )
    estimated_chars = 0
    total_files = len(selected_files)
    files_with_content_flag = [
        *((item, True) for item in content_files),
        *((item, False) for item in filename_only_files),
    ]
    cost_workers = min(FILE_EXTRACTION_MAX_WORKERS, len(files_with_content_flag))
    completed_files = 0
    if cost_workers:
        executor = ThreadPoolExecutor(max_workers=cost_workers)
        try:
            for batch_start in range(0, len(files_with_content_flag), cost_workers):
                _raise_if_cancelled(cancel_check)
                batch = files_with_content_flag[
                    batch_start : batch_start + cost_workers
                ]
                futures = [
                    executor.submit(
                        _estimate_file_chars,
                        file,
                        absolute_path,
                        include_content,
                    )
                    for (file, absolute_path), include_content in batch
                ]
                deadline = time.monotonic() + COST_ESTIMATION_TIMEOUT_SECONDS
                batch_timed_out = False
                for future, ((file, absolute_path), _include_content) in zip(
                    futures, batch
                ):
                    _raise_if_cancelled(cancel_check)
                    try:
                        file_chars = future.result(
                            timeout=max(0.0, deadline - time.monotonic())
                        )
                    except TimeoutError:
                        logging.warning(
                            "비용 추산이 5초를 초과해 평균값으로 대체합니다: %s",
                            absolute_path,
                        )
                        file_chars = 2_000
                        batch_timed_out = True
                        future.cancel()
                    except Exception:
                        logging.warning(
                            "비용 추산 중 오류가 발생해 평균값으로 대체합니다: %s",
                            absolute_path,
                        )
                        file_chars = 2_000
                    estimated_chars += file_chars
                    completed_files += 1
                    if progress_callback is not None:
                        progress_callback("cost", completed_files, total_files)
                if batch_timed_out:
                    # Running Python threads cannot be killed safely. Abandon
                    # only the affected pool and continue with fresh capacity.
                    executor.shutdown(wait=False, cancel_futures=True)
                    executor = ThreadPoolExecutor(max_workers=cost_workers)
        finally:
            executor.shutdown(wait=False, cancel_futures=True)
    for memo in analysis_result.memos:
        _raise_if_cancelled(cancel_check)
        estimated_chars += len(f"{memo.title}\n{memo.content}".strip())
    for record in _get_email_text_records(parsed_emails):
        _raise_if_cancelled(cancel_check)
        estimated_chars += len(preprocess_text_for_rag(record["text"], "email"))
    for record in _get_email_attachment_text_records(parsed_emails):
        _raise_if_cancelled(cancel_check)
        estimated_chars += len(preprocess_text_for_rag(record["text"], "document"))
    for path, raw_text in _get_kakao_text_records(kakao_file_paths):
        _raise_if_cancelled(cancel_check)
        estimated_chars += len(preprocess_text_for_rag(raw_text, "kakao"))
    # Korean business documents have measured close to one embedding token per
    # extracted character.  Reserve conservatively so a package cannot start
    # with only the old fixed 40-credit allowance.
    estimated_tokens = max(1, math.ceil(estimated_chars * 1.0)) if estimated_chars else 0
    estimated_cost_krw = (
        math.ceil((estimated_tokens / 1000) * embedding_unit_cost_per_1k)
        if embedding_unit_cost_per_1k is not None
        and embedding_unit_cost_per_1k > 0
        else 0
    )
    estimated_size_bytes = sum(file.size_bytes for file, _path in content_files)
    email_source_paths = {
        str(record.get("source_file") or "").split("::", 1)[0]
        for record in (parsed_emails or [])
        if isinstance(record, dict)
    }
    for path in [*email_source_paths, *(kakao_file_paths or [])]:
        try:
            estimated_size_bytes += Path(path).stat().st_size
        except OSError:
            pass
    return {
        "file_count": len(selected_files),
        "content_file_count": len(content_files),
        "estimated_tokens": estimated_tokens,
        "estimated_cost_krw": estimated_cost_krw,
        "embedding_unit_cost_per_1k": embedding_unit_cost_per_1k or 0,
        "estimated_size_bytes": estimated_size_bytes,
    }


def _get_rag_package_zip_path(output_path: Path) -> Path:
    if output_path.suffix.casefold() == ".zip":
        return output_path
    return output_path.with_suffix(".zip")


def _zip_package_folder(
    package_path: Path,
    zip_path: Path,
    cancel_check: Callable[[], bool] | None = None,
) -> None:
    with zipfile.ZipFile(zip_path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for file_path in sorted(path for path in package_path.rglob("*") if path.is_file()):
            _raise_if_cancelled(cancel_check)
            archive.write(file_path, file_path.relative_to(package_path))


def _file_extraction_timeout(file_name: str) -> int:
    if Path(file_name).suffix.lower() in {".xlsx", ".xls", ".hwp", ".hwpx"}:
        return SLOW_FILE_EXTRACTION_TIMEOUT_SECONDS
    return DEFAULT_FILE_EXTRACTION_TIMEOUT_SECONDS


def _terminate_process_executor(executor: ProcessPoolExecutor) -> None:
    terminate_process_executor(executor)


def _extract_text_in_process(file_path: str) -> str:
    return extract_text_for_rag(file_path)


def _extract_file_with_hard_timeout(
    file_path: Path,
    timeout_seconds: float,
    extractor: Callable[[str], str] = _extract_text_in_process,
) -> str:
    executor = ProcessPoolExecutor(max_workers=1)
    future = executor.submit(extractor, str(file_path))
    try:
        text = future.result(timeout=timeout_seconds)
    except TimeoutError:
        future.cancel()
        _terminate_process_executor(executor)
        raise
    except BaseException:
        _terminate_process_executor(executor)
        raise
    else:
        executor.shutdown(wait=True)
        return text


def _build_file_chunk_records_with_timeout(
    analyzed_file: AnalyzedFile,
    absolute_path: Path,
    include_content: bool = True,
) -> tuple[list[dict], list[dict], bool]:
    if not include_content:
        return ([
            _build_not_in_whitelist_placeholder_chunk_record(
                analyzed_file.relative_path,
                analyzed_file.file_name,
                analyzed_file.relative_path,
                analyzed_file.modified_at,
            )
        ], [], False)
    if not _is_rag_text_extraction_file(analyzed_file.file_name):
        records, exclusions = _build_file_chunk_records(analyzed_file, absolute_path)
        return records, exclusions, False
    try:
        raw_text = _extract_file_with_hard_timeout(
            absolute_path, _file_extraction_timeout(analyzed_file.file_name)
        )
    except TimeoutError:
        logging.warning("파일 추출 시간 초과로 건너뜁니다: %s", absolute_path)
        return (
            [_build_timeout_placeholder_chunk_record(
                analyzed_file.relative_path,
                analyzed_file.file_name,
                analyzed_file.relative_path,
                analyzed_file.modified_at,
            )],
            [_build_excluded_file_record(analyzed_file, "extract_timeout")],
            True,
        )
    text = preprocess_text_for_rag(raw_text, _infer_source_type(analyzed_file))
    return (*_build_file_chunk_records_from_text(analyzed_file, text), False)


def _timeout_file_result(
    analyzed_file: AnalyzedFile,
) -> tuple[list[dict], list[dict], bool]:
    return (
        [_build_timeout_placeholder_chunk_record(
            analyzed_file.relative_path,
            analyzed_file.file_name,
            analyzed_file.relative_path,
            analyzed_file.modified_at,
        )],
        [_build_excluded_file_record(analyzed_file, "extract_timeout")],
        True,
    )


def _finish_extracted_file(
    analyzed_file: AnalyzedFile,
    raw_text: str,
) -> tuple[list[dict], list[dict], bool]:
    text = preprocess_text_for_rag(raw_text, _infer_source_type(analyzed_file))
    return (*_build_file_chunk_records_from_text(analyzed_file, text), False)


def _process_pending_files_with_shared_pool(
    pending: list[tuple[int, AnalyzedFile, Path, bool]],
    selected_files: list[tuple[AnalyzedFile, Path]],
    identities: list[str],
    file_results: list[tuple[list[dict], list[dict]] | None],
    checkpoint_path: Path,
    timed_out_files: list[str],
    completed_files: int,
    total_files: int,
    progress_callback: Callable[[str, int, int], None] | None,
    cancel_check: Callable[[], bool] | None,
) -> int:
    extraction_queue: deque[tuple[int, AnalyzedFile, Path]] = deque()

    def complete(
        index: int,
        records: list[dict],
        exclusions: list[dict],
        timed_out: bool,
    ) -> None:
        nonlocal completed_files
        file_results[index] = (records, exclusions)
        timed_out_file = ""
        if timed_out:
            timed_out_file = selected_files[index][0].relative_path
            if timed_out_file not in timed_out_files:
                timed_out_files.append(timed_out_file)
        _append_checkpoint_result(
            checkpoint_path,
            identities[index],
            records,
            exclusions,
            timed_out_file,
        )
        completed_files += 1
        if progress_callback is not None:
            progress_callback("files", completed_files, total_files)
        _raise_if_cancelled(cancel_check)

    # Filename-only records do not need a process. Completing them in the
    # parent keeps all worker slots available for actual document extraction.
    for index, analyzed_file, absolute_path, include_content in pending:
        _raise_if_cancelled(cancel_check)
        if include_content:
            extraction_queue.append((index, analyzed_file, absolute_path))
            continue
        records, exclusions, timed_out = _build_file_chunk_records_with_timeout(
            analyzed_file,
            absolute_path,
            False,
        )
        complete(index, records, exclusions, timed_out)

    if not extraction_queue:
        return completed_files

    worker_count = min(FILE_EXTRACTION_MAX_WORKERS, len(extraction_queue))
    executor = ProcessPoolExecutor(max_workers=worker_count)
    in_flight: dict = {}
    try:
        while extraction_queue or in_flight:
            _raise_if_cancelled(cancel_check)
            while extraction_queue and len(in_flight) < worker_count:
                task = extraction_queue.popleft()
                _index, _analyzed_file, absolute_path = task
                future = executor.submit(_extract_text_in_process, str(absolute_path))
                in_flight[future] = (task, time.monotonic())

            next_deadline = min(
                submitted_at + _file_extraction_timeout(task[1].file_name)
                for task, submitted_at in in_flight.values()
            )
            completed, _ = wait(
                in_flight,
                timeout=min(0.25, max(0.0, next_deadline - time.monotonic())),
                return_when=FIRST_COMPLETED,
            )
            for future in completed:
                (index, analyzed_file, _absolute_path), _submitted_at = (
                    in_flight.pop(future)
                )
                try:
                    raw_text = future.result()
                except BaseException:
                    logging.exception("파일 추출 프로세스 오류: %s", analyzed_file.relative_path)
                    records, exclusions, timed_out = _build_file_chunk_records(
                        analyzed_file, selected_files[index][1]
                    ) + (False,)
                else:
                    records, exclusions, timed_out = _finish_extracted_file(
                        analyzed_file, raw_text
                    )
                complete(index, records, exclusions, timed_out)

            now = time.monotonic()
            expired = [
                future
                for future, (task, submitted_at) in in_flight.items()
                if now - submitted_at >= _file_extraction_timeout(task[1].file_name)
            ]
            if not expired:
                continue

            expired_set = set(expired)
            retry_tasks: list[tuple[int, AnalyzedFile, Path]] = []
            for future, (task, _submitted_at) in list(in_flight.items()):
                index, analyzed_file, absolute_path = task
                if future in expired_set:
                    logging.warning(
                        "파일 추출 시간 초과로 건너뜁니다: %s", absolute_path
                    )
                    complete(index, *_timeout_file_result(analyzed_file))
                else:
                    retry_tasks.append(task)
            in_flight.clear()
            terminate_process_executor(executor)
            extraction_queue.extendleft(reversed(retry_tasks))
            if extraction_queue:
                executor = ProcessPoolExecutor(max_workers=worker_count)

        executor.shutdown(wait=True)
    except BaseException:
        terminate_process_executor(executor)
        raise
    return completed_files


def _checkpoint_identity(analyzed_file: AnalyzedFile, absolute_path: Path) -> str:
    stat = absolute_path.stat()
    value = f"{absolute_path.resolve()}|{stat.st_size}|{stat.st_mtime_ns}"
    return hashlib.sha256(value.encode("utf-8", errors="surrogatepass")).hexdigest()


def _checkpoint_path(folder_paths: list[str]) -> Path:
    roots = "|".join(sorted(str(Path(path).resolve()) for path in folder_paths))
    key = hashlib.sha256(roots.encode("utf-8", errors="surrogatepass")).hexdigest()[:20]
    return Path("config") / "rag_package_checkpoints" / f"{key}.jsonl"


def _load_checkpoint(path: Path, signature: str) -> dict[str, Any]:
    candidate_path = path
    if not candidate_path.exists() and path.suffix == ".jsonl":
        legacy_path = path.with_suffix(".json")
        if legacy_path.exists():
            candidate_path = legacy_path

    try:
        raw_text = candidate_path.read_text(encoding="utf-8")
    except OSError:
        raw_text = ""

    try:
        data = json.loads(raw_text)
        if data.get("signature") == signature and isinstance(data.get("results"), dict):
            return data
    except (ValueError, TypeError, AttributeError):
        pass

    if raw_text:
        results: dict[str, Any] = {}
        timed_out_files: list[str] = []
        metadata_seen = False
        for raw_line in raw_text.splitlines():
            if not raw_line.strip():
                continue
            try:
                record = json.loads(raw_line)
            except (ValueError, TypeError):
                # A process may be stopped while writing the last line. All
                # earlier complete JSONL records remain resumable.
                continue
            if record.get("type") == "meta":
                if record.get("signature") != signature:
                    return {"signature": signature, "results": {}, "timed_out_files": []}
                metadata_seen = True
                timed_out_files.extend(record.get("timed_out_files") or [])
                continue
            if record.get("type") != "result" or not metadata_seen:
                continue
            identity = record.get("identity")
            value = record.get("value")
            if isinstance(identity, str) and isinstance(value, list) and len(value) == 2:
                results[identity] = value
                timed_out_file = record.get("timed_out_file")
                if timed_out_file:
                    timed_out_files.append(str(timed_out_file))
        if metadata_seen:
            return {
                "signature": signature,
                "results": results,
                "timed_out_files": list(dict.fromkeys(timed_out_files)),
            }
    return {"signature": signature, "results": {}, "timed_out_files": []}


def _save_checkpoint(path: Path, data: dict[str, Any]) -> None:
    """Write a complete JSONL checkpoint for initialization or migration."""
    path.parent.mkdir(parents=True, exist_ok=True)
    temporary_path = path.with_suffix(".tmp")
    timed_out_files = list(data.get("timed_out_files") or [])
    lines = [json.dumps({
        "type": "meta",
        "signature": data.get("signature", ""),
        "timed_out_files": timed_out_files,
    }, ensure_ascii=False, separators=(",", ":"))]
    for identity, value in (data.get("results") or {}).items():
        lines.append(json.dumps({
            "type": "result",
            "identity": identity,
            "value": value,
        }, ensure_ascii=False, separators=(",", ":")))
    temporary_path.write_text("\n".join(lines) + "\n", encoding="utf-8")
    os.replace(temporary_path, path)


def _is_current_jsonl_checkpoint(path: Path, signature: str) -> bool:
    try:
        with path.open("r", encoding="utf-8") as checkpoint_file:
            first_line = checkpoint_file.readline()
        metadata = json.loads(first_line)
        return (
            metadata.get("type") == "meta"
            and metadata.get("signature") == signature
        )
    except (OSError, ValueError, TypeError, AttributeError):
        return False


def _append_checkpoint_result(
    path: Path,
    identity: str,
    records: list[dict],
    exclusions: list[dict],
    timed_out_file: str = "",
) -> None:
    line = json.dumps({
        "type": "result",
        "identity": identity,
        "value": [records, exclusions],
        "timed_out_file": timed_out_file,
    }, ensure_ascii=False, separators=(",", ":"))
    with path.open("a", encoding="utf-8") as checkpoint_file:
        checkpoint_file.write(line + "\n")


def _build_checkpoint_signature(
    identities: list[str],
    selected_extensions: set[str],
    max_file_size_bytes: int | None,
    extension_size_limits: dict[str, int | None] | None,
) -> str:
    signature_input = "|".join([
        *identities,
        "selected=" + ",".join(sorted(selected_extensions)),
        "max_size=" + (
            str(max_file_size_bytes) if max_file_size_bytes is not None else "unlimited"
        ),
        "extension_limits=" + ",".join(
            f"{extension}:{'unlimited' if limit is None else limit}"
            for extension, limit in sorted((extension_size_limits or {}).items())
        ),
    ])
    return hashlib.sha256(signature_input.encode("ascii")).hexdigest()


def build_and_save_rag_package(
    analysis_result: AnalysisResult,
    folder_paths: list[str],
    api_key: str,
    output_path: str,
    parsed_emails: list[dict] | None = None,
    kakao_file_paths: list[str] | None = None,
    progress_callback: Callable[[str, int, int], None] | None = None,
    cancel_check: Callable[[], bool] | None = None,
    selected_extensions: set[str] | None = None,
    max_file_size_bytes: int | None = None,
    extension_size_limits: dict[str, int | None] | None = None,
) -> dict[str, Any]:
    _raise_if_cancelled(cancel_check)
    selected_files, excluded_files = _collect_rag_package_candidate_files(
        analysis_result,
        folder_paths,
        cancel_check,
    )
    normalized_selected_extensions = (
        {extension.lower() for extension in selected_extensions}
        if selected_extensions is not None
        else set(RAG_TEXT_EXTRACTION_EXTENSIONS)
    )
    content_files, _filename_only_files = filter_files_by_selected_extensions(
        selected_files,
        normalized_selected_extensions,
        max_file_size_bytes,
        extension_size_limits,
    )
    content_paths = {str(absolute_path) for _file, absolute_path in content_files}

    chunk_records: list[dict] = []
    source_map: dict[str, dict[str, Any]] = {}
    total_files = len(selected_files)
    file_results: list[tuple[list[dict], list[dict]] | None] = [None] * total_files
    identities = [_checkpoint_identity(file, path) for file, path in selected_files]
    signature = _build_checkpoint_signature(
        identities,
        normalized_selected_extensions,
        max_file_size_bytes,
        extension_size_limits,
    )
    checkpoint_path = _checkpoint_path(folder_paths)
    checkpoint = _load_checkpoint(checkpoint_path, signature)
    checkpoint_results = checkpoint["results"]
    timed_out_files = list(checkpoint.get("timed_out_files", []))
    if not _is_current_jsonl_checkpoint(checkpoint_path, signature):
        _save_checkpoint(checkpoint_path, checkpoint)
    pending: list[tuple[int, AnalyzedFile, Path, bool]] = []
    for index, ((analyzed_file, absolute_path), identity) in enumerate(zip(selected_files, identities)):
        saved_result = checkpoint_results.get(identity)
        if isinstance(saved_result, list) and len(saved_result) == 2:
            file_results[index] = (saved_result[0], saved_result[1])
        else:
            pending.append((
                index,
                analyzed_file,
                absolute_path,
                str(absolute_path) in content_paths,
            ))
    completed_files = total_files - len(pending)
    if completed_files and progress_callback is not None:
        progress_callback("resume", completed_files, total_files)
    if pending:
        completed_files = _process_pending_files_with_shared_pool(
            pending,
            selected_files,
            identities,
            file_results,
            checkpoint_path,
            timed_out_files,
            completed_files,
            total_files,
            progress_callback,
            cancel_check,
        )

    for result_records in file_results:
        _raise_if_cancelled(cancel_check)
        if result_records is None:
            continue
        records, exclusions = result_records
        chunk_records.extend(records)
        excluded_files.extend(exclusions)

    for memo_index, memo in enumerate(analysis_result.memos):
        _raise_if_cancelled(cancel_check)
        memo_text = f"{memo.title}\n{memo.content}".strip()
        for chunk_index, chunk in enumerate(split_into_chunks(memo_text)):
            source_file = f"memo:{memo_index}"
            chunk_records.append(
                {
                    "chunk_text": chunk,
                    "source_file": source_file,
                    "chunk_index": chunk_index,
                    "metadata": {
                        "source_path": source_file,
                        "file_name": f"메모: {memo.title}",
                        "modified_at": memo.updatedat,
                    },
                }
            )

    for record in _get_email_text_records(parsed_emails):
        _raise_if_cancelled(cancel_check)
        text = preprocess_text_for_rag(record["text"], "email")
        for chunk_index, chunk in enumerate(split_into_chunks(text)):
            source_file = record["source_file"]
            chunk_records.append(
                {
                    "chunk_text": chunk,
                    "source_file": source_file,
                    "chunk_index": chunk_index,
                    "metadata": {
                        "source_path": source_file,
                        "file_name": f"메일: {record['subject']}",
                        "modified_at": record["date"],
                    },
                }
            )

    prepared_attachment_text = {
        _attachment_cache_key(record): record.get("text", "")
        for record in _get_email_attachment_text_records(parsed_emails)
        if _is_supported_email_attachment(record)
    }
    for record in _get_email_attachment_records(parsed_emails):
        _raise_if_cancelled(cancel_check)
        source_file = f"{record['source_file']}::attachment::{record['filename']}"
        attachment_display_name = (
            f"{record['email_subject']}의 첨부파일: {record['filename']}"
        )
        if not _is_supported_email_attachment(record):
            chunk_records.append(
                _build_not_in_whitelist_placeholder_chunk_record(
                    source_file,
                    attachment_display_name,
                    source_file,
                    record["date"],
                )
            )
            continue

        text = preprocess_text_for_rag(
            str(prepared_attachment_text.get(_attachment_cache_key(record), "")),
            "document",
        )
        if not text:
            excluded_files.append(
                {
                    "file_name": record["filename"],
                    "relative_path": (
                        f"출처: {record['email_subject']}의 첨부파일 "
                        f"{record['filename']}"
                    ),
                    "modified_at": record["date"],
                    "reason": "email_attachment_extract_failed",
                }
            )
            chunk_records.append(
                _build_unsupported_placeholder_chunk_record(
                    source_file,
                    attachment_display_name,
                    source_file,
                    record["date"],
                )
            )
            continue

        for chunk_index, chunk in enumerate(
            _split_text_for_file(text, record["filename"])
        ):
            chunk_records.append(
                {
                    "chunk_text": chunk,
                    "source_file": source_file,
                    "chunk_index": chunk_index,
                    "metadata": {
                        "source_path": source_file,
                        "file_name": (
                            f"{record['email_subject']} - 첨부파일: "
                            f"{record['filename']}"
                        ),
                        "modified_at": record["date"],
                    },
                }
            )

    for external_path, raw_text in _get_kakao_text_records(kakao_file_paths):
        _raise_if_cancelled(cancel_check)
        text = preprocess_text_for_rag(raw_text, "kakao")
        for chunk_index, chunk in enumerate(split_into_chunks(text)):
            chunk_records.append(
                {
                    "chunk_text": chunk,
                    "source_file": external_path,
                    "chunk_index": chunk_index,
                    "metadata": {
                        "source_path": external_path,
                        "file_name": Path(external_path).name,
                        "modified_at": "",
                    },
                }
            )

    embedded, embedding_failures, embedding_tokens = _embed_chunk_records(
        chunk_records,
        api_key,
        progress_callback,
        cancel_check,
    )
    for index, item in enumerate(embedded):
        _raise_if_cancelled(cancel_check)
        chunk_id = f"chunk-{index:06d}"
        item["chunk_id"] = chunk_id
        source_map[chunk_id] = item.pop("metadata", {})

    saved_path = save_rag_package(
        embedded,
        analysis_result,
        output_path,
        api_key,
        source_map=source_map,
        excluded_files=excluded_files,
        embedding_failures=embedding_failures,
        cancel_check=cancel_check,
    )
    try:
        checkpoint_path.unlink(missing_ok=True)
        if checkpoint_path.suffix == ".jsonl":
            checkpoint_path.with_suffix(".json").unlink(missing_ok=True)
    except OSError:
        logging.warning("완료된 체크포인트를 삭제하지 못했습니다: %s", checkpoint_path)
    return {
        "content_file_count": len(content_files),
        "filename_only_file_count": len(_filename_only_files),
        "embedding_failed_chunk_count": len(embedding_failures),
        "embedding_tokens": embedding_tokens,
        "saved_path": saved_path,
        "timed_out_file_count": len(timed_out_files),
        "timed_out_files": timed_out_files,
    }


def _estimate_file_chars(
    file: AnalyzedFile,
    absolute_path: Path,
    include_content: bool,
) -> int:
    if not include_content:
        return len(_build_not_in_whitelist_placeholder_text(
            file.file_name, file.relative_path, file.modified_at
        ))
    suffix = Path(file.file_name).suffix.lower()
    if suffix in {".xlsx", ".xls"}:
        return estimate_excel_cost_lightweight(absolute_path)
    if suffix in {".hwp", ".hwpx"}:
        return estimate_hwp_cost_lightweight(absolute_path)
    return len(_get_preprocessed_text_for_file(file, absolute_path))


def _build_timeout_placeholder_chunk_record(
    source_file: str,
    file_name: str,
    source_path: str,
    modified_at: str,
) -> dict:
    return {
        "chunk_text": (
            f"[파일명: {file_name}] 이 파일은 처리 시간 초과로 건너뛰었습니다. "
            f"원본 경로: {source_path}, 수정일시: {modified_at}"
        ),
        "source_file": source_file,
        "chunk_index": 0,
        "metadata": {
            "source_path": source_path,
            "file_name": file_name,
            "modified_at": modified_at,
            "extraction_status": "extract_timeout",
        },
    }


def estimate_excel_cost_lightweight(file_path: Path) -> int:
    """Estimate spreadsheet text size without running the full RAG extractor.

    This intentionally uses only sheet dimensions and file size. It does not
    inspect cell values, formulas, merged ranges, or number formats; accurate
    extraction remains exclusively in the actual embedding stage.
    """
    max_rows = MAX_EXCEL_EMBED_ROWS + 1  # include the header row
    max_columns = MAX_EXCEL_EMBED_COLUMNS
    estimated_cells = 0
    try:
        if file_path.suffix.lower() == ".xlsx":
            from openpyxl import load_workbook

            with warnings.catch_warnings():
                warnings.simplefilter("ignore", UserWarning)
                workbook = load_workbook(file_path, read_only=True, data_only=True)
            try:
                for worksheet in workbook.worksheets:
                    estimated_cells += min(worksheet.max_row, max_rows) * min(
                        worksheet.max_column, max_columns
                    )
            finally:
                workbook.close()
        else:
            import xlrd

            workbook = xlrd.open_workbook(file_path, on_demand=True)
            try:
                for worksheet in workbook.sheets():
                    estimated_cells += min(worksheet.nrows, max_rows) * min(
                        worksheet.ncols, max_columns
                    )
            finally:
                workbook.release_resources()
    except Exception:
        return 0

    # A conservative average cell width plus a small file-size component gives
    # a useful order-of-magnitude estimate without reading cell contents.
    file_size_chars = min(max(0, int(os.path.getsize(file_path) / 2)), 20_000)
    return estimated_cells * 12 + file_size_chars


def estimate_hwp_cost_lightweight(file_path: Path) -> int:
    """Estimate HWP/HWPX text size from container metadata only."""
    try:
        size = max(0, int(os.path.getsize(file_path)))
        if file_path.suffix.lower() == ".hwp":
            import olefile

            with olefile.OleFileIO(str(file_path)) as ole:
                body_sizes = [ole.get_size(name) for name in ole.listdir()
                              if name and name[0].lower() == "bodytext"]
            return max(1, int(sum(body_sizes) * 0.35))
        with zipfile.ZipFile(file_path) as archive:
            xml_size = sum(info.file_size for info in archive.infolist()
                           if info.filename.lower().endswith(".xml"))
        return max(1, int(xml_size * 0.12), int(size * 0.5))
    except Exception:
        return max(1, int(size * 0.5))


def _build_file_chunk_records(
    analyzed_file: AnalyzedFile,
    absolute_path: Path,
) -> tuple[list[dict], list[dict]]:
    exclusions: list[dict] = []
    records: list[dict] = []
    if not _is_rag_text_extraction_file(analyzed_file.file_name):
        records.append(
            _build_not_in_whitelist_placeholder_chunk_record(
                analyzed_file.relative_path,
                analyzed_file.file_name,
                analyzed_file.relative_path,
                analyzed_file.modified_at,
            )
        )
        return records, exclusions

    text = _get_preprocessed_text_for_file(analyzed_file, absolute_path)
    if _is_pdf_extract_timeout(absolute_path):
        exclusions.append(_build_excluded_file_record(analyzed_file, "pdf_extract_timeout"))
        return records, exclusions
    return _build_file_chunk_records_from_text(analyzed_file, text)


def _build_file_chunk_records_from_text(
    analyzed_file: AnalyzedFile,
    text: str,
) -> tuple[list[dict], list[dict]]:
    exclusions: list[dict] = []
    records: list[dict] = []
    if not text:
        exclusions.append(_build_excluded_file_record(analyzed_file, "unsupported_format"))
        records.append(
            _build_unsupported_placeholder_chunk_record(
                analyzed_file.relative_path,
                analyzed_file.file_name,
                analyzed_file.relative_path,
                analyzed_file.modified_at,
            )
        )
        return records, exclusions

    for chunk_index, chunk in enumerate(
        _split_text_for_file(text, analyzed_file.file_name)
    ):
        records.append(
            {
                "chunk_text": chunk,
                "source_file": analyzed_file.relative_path,
                "chunk_index": chunk_index,
                "metadata": {
                    "source_path": analyzed_file.relative_path,
                    "file_name": analyzed_file.file_name,
                    "modified_at": analyzed_file.modified_at,
                },
            }
        )
    return records, exclusions


def _remove_email_reply_chain(text: str) -> str:
    patterns = [
        r"(?im)^-----Original Message-----.*\Z",
        r"(?im)^>.*(?:\n>.*)*",
        r"(?im)^On .+ wrote:\s*.*\Z",
    ]
    result = text
    for pattern in patterns:
        result = re.sub(pattern, "", result, flags=re.DOTALL)
    return result


def _remove_lines_with_keywords(text: str, keywords: tuple[str, ...]) -> str:
    return "\n".join(
        line
        for line in text.splitlines()
        if not any(keyword in line.casefold() for keyword in keywords)
    )


def _looks_like_html(text: str) -> bool:
    lowered = text[:2000].casefold()
    return (
        "<html" in lowered
        or "<!doctype html" in lowered
        or bool(re.search(r"<(?:body|head|div|span|p|table|br)\b", lowered))
    )


def _strip_html_to_text(text: str) -> str:
    without_scripts = re.sub(
        r"(?is)<(script|style)\b[^>]*>.*?</\1>",
        " ",
        text,
    )
    with_line_breaks = re.sub(
        r"(?i)<\s*(br|/p|/div|/tr|/li|/h[1-6])\s*/?\s*>",
        "\n",
        without_scripts,
    )
    without_tags = re.sub(r"(?s)<[^>]+>", " ", with_line_breaks)
    decoded = unescape(without_tags)
    lines = [
        re.sub(r"[ \t\r\f\v]+", " ", line).strip()
        for line in decoded.splitlines()
    ]
    compact_lines = [line for line in lines if line]
    return "\n".join(compact_lines)


def _remove_email_signature(text: str) -> str:
    lines = text.splitlines()
    for index, line in enumerate(lines):
        stripped = line.strip()
        if re.fullmatch(r"[-_=]{2,}", stripped):
            signature = "\n".join(lines[index + 1 : index + 8]).casefold()
            hits = sum(1 for keyword in EMAIL_SIGNATURE_KEYWORDS if keyword in signature)
            if hits >= 2:
                return "\n".join(lines[:index]).strip()
    return text


def _split_long_text(text: str, chunk_size: int, overlap: int) -> list[str]:
    if len(text) <= chunk_size:
        return [text]
    parts = []
    start = 0
    step = max(chunk_size - overlap, 1)
    while start < len(text):
        parts.append(text[start : start + chunk_size].strip())
        start += step
    return [part for part in parts if part]


def _tail_overlap(text: str, overlap: int) -> str:
    if overlap <= 0 or len(text) <= overlap:
        return text if len(text) <= overlap else ""
    return text[-overlap:].strip()


def _create_embedding_with_retry(client: Any, batch: list[str]) -> Any:
    last_exception: Exception | None = None
    for attempt in range(EMBEDDING_RETRY_ATTEMPTS):
        try:
            return client.embeddings.create(model=EMBEDDING_MODEL, input=batch)
        except Exception as exc:
            last_exception = exc
            if attempt == EMBEDDING_RETRY_ATTEMPTS - 1:
                break
            time.sleep(2 ** (attempt + 1))
    detail = (
        f"{type(last_exception).__name__}: {last_exception}"
        if last_exception is not None
        else "unknown error"
    )
    raise RuntimeError(f"임베딩 API 호출에 실패했습니다: {detail}") from last_exception


def _embed_chunk_records(
    chunk_records: list[dict],
    api_key: str,
    progress_callback: Callable[[str, int, int], None] | None = None,
    cancel_check: Callable[[], bool] | None = None,
) -> tuple[list[dict], list[dict], int]:
    texts = [record["chunk_text"] for record in chunk_records]
    embedded, failures, embedding_tokens = _embed_chunks_with_failures(
        texts,
        api_key,
        progress_callback,
        cancel_check,
    )
    for item in embedded:
        original_index = int(item["chunk_index"])
        record = chunk_records[original_index]
        item["source_file"] = record["source_file"]
        item["chunk_index"] = record["chunk_index"]
        item["metadata"] = record["metadata"]

    failure_records: list[dict] = []
    for failure in failures:
        original_index = int(failure["chunk_index"])
        record = chunk_records[original_index]
        failure_records.append(
            {
                "source_file": record["source_file"],
                "chunk_index": record["chunk_index"],
                "metadata": record["metadata"],
                "chunk_text_preview": failure["chunk_text_preview"],
                "batch_index": failure["batch_index"],
                "error": failure["error"],
            }
        )
    return embedded, failure_records, embedding_tokens


def _get_preprocessed_text_for_file(file: AnalyzedFile, absolute_path: Path) -> str:
    text = _get_cached_text_for_path(str(absolute_path))
    return preprocess_text_for_rag(text, _infer_source_type(file))


def _get_cached_text_for_path(path: str) -> str:
    cache_key = str(Path(path))
    if cache_key not in _EXTRACTED_TEXT_CACHE:
        _EXTRACTED_TEXT_CACHE[cache_key] = extract_text_for_rag(cache_key)
    return _EXTRACTED_TEXT_CACHE[cache_key]


def _is_rag_text_extraction_file(file_name: str) -> bool:
    return Path(file_name).suffix.lower() in RAG_TEXT_EXTRACTION_EXTENSIONS


def _is_pdf_extract_timeout(path: Path) -> bool:
    return str(path) in _PDF_TEXT_TIMEOUT_PATHS


def _get_email_text_records(
    parsed_emails: list[dict] | None,
) -> list[dict]:
    if not parsed_emails:
        return []

    records: list[dict] = []
    seen: set[str] = set()
    for email in parsed_emails:
        try:
            source_file = str(email.get("source_file", ""))
            if source_file in seen:
                continue
            seen.add(source_file)
            subject = str(email.get("subject") or "(제목 없음)")
            sender = str(email.get("sender") or "")
            date = str(email.get("date") or "")
            body = str(email.get("body") or "")
            records.append(
                {
                    "source_file": source_file,
                    "subject": subject,
                    "date": date,
                    "text": f"발신자: {sender}\n발송일시: {date}\n제목: {subject}\n\n{body}".strip(),
                }
            )
        except Exception:
            continue
    return records


def _get_email_attachment_text_records(
    parsed_emails: list[dict] | None,
) -> list[dict]:
    records: list[dict] = []
    supported_records: list[dict] = []
    for record in _get_email_attachment_records(parsed_emails):
        if not _is_supported_email_attachment(record):
            attachment_display_name = (
                f"{record['email_subject']}의 첨부파일: {record['filename']}"
            )
            source_file = f"{record['source_file']}::attachment::{record['filename']}"
            records.append(
                {
                    **record,
                    "text": _build_not_in_whitelist_placeholder_text(
                        attachment_display_name,
                        source_file,
                        record["date"],
                    ),
                }
            )
            continue
        cache_key = _attachment_cache_key(record)
        if cache_key in _ATTACHMENT_TEXT_CACHE:
            cached_text = _ATTACHMENT_TEXT_CACHE[cache_key]
            if cached_text:
                records.append({**record, "text": cached_text})
            continue
        supported_records.append(record)

    outcomes = run_process_items_with_timeout(
        supported_records,
        _extract_attachment_text_worker,
        timeout_seconds=SLOW_FILE_EXTRACTION_TIMEOUT_SECONDS,
        max_workers=FILE_EXTRACTION_MAX_WORKERS,
    )
    for record, (status, text) in zip(supported_records, outcomes):
        if status != "ok" or not text:
            continue
        _ATTACHMENT_TEXT_CACHE[_attachment_cache_key(record)] = text
        records.append({**record, "text": text})
    return records


def _extract_attachment_text_worker(record: dict) -> str:
    return _extract_attachment_text(record)


def _attachment_cache_key(record: dict) -> str:
    content_bytes = record.get("content_bytes") or b""
    extension = Path(record.get("filename", "")).suffix.lower()
    return hashlib.sha256(
        extension.encode("utf-8") + b"\0" + content_bytes
    ).hexdigest()


def _get_email_attachment_records(
    parsed_emails: list[dict] | None,
) -> list[dict]:
    if not parsed_emails:
        return []

    records: list[dict] = []
    seen: set[str] = set()
    for email in parsed_emails:
        try:
            source_file = str(email.get("source_file", ""))
            email_subject = str(email.get("subject") or "(제목 없음)")
            date = str(email.get("date") or "")
            for attachment in email.get("attachments", []) or []:
                filename = str(attachment.get("filename") or "attachment")
                content_bytes = attachment.get("content_bytes") or b""
                if not isinstance(content_bytes, bytes):
                    continue
                key = f"{source_file}\0{filename}\0{hashlib.sha256(content_bytes).hexdigest()}"
                if key in seen:
                    continue
                seen.add(key)
                records.append(
                    {
                        "source_file": source_file,
                        "email_subject": email_subject,
                        "date": date,
                        "filename": filename,
                        "content_bytes": content_bytes,
                        "size": int(attachment.get("size") or len(content_bytes)),
                    }
                )
        except Exception:
            continue
    return records


def _is_supported_email_attachment(record: dict) -> bool:
    return Path(record.get("filename", "")).suffix.lower() in SUPPORTED_EMAIL_ATTACHMENT_EXTENSIONS


def _extract_attachment_text(record: dict) -> str:
    content_bytes = record.get("content_bytes") or b""
    if not isinstance(content_bytes, bytes) or not content_bytes:
        return ""

    extension = Path(record.get("filename", "")).suffix.lower()
    cache_key = _attachment_cache_key(record)
    if cache_key in _ATTACHMENT_TEXT_CACHE:
        return _ATTACHMENT_TEXT_CACHE[cache_key]

    temp_path = ""
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=extension) as temp_file:
            temp_file.write(content_bytes)
            temp_path = temp_file.name
        text = extract_text_for_rag(temp_path)
    except Exception:
        text = ""
    finally:
        if temp_path:
            try:
                os.unlink(temp_path)
            except OSError:
                pass

    _ATTACHMENT_TEXT_CACHE[cache_key] = text
    return text


def _get_kakao_text_paths(kakao_file_paths: list[str] | None) -> list[str]:
    paths: list[str] = []
    seen: set[str] = set()
    for path in kakao_file_paths or []:
        if path in seen or not Path(path).is_file():
            continue
        seen.add(path)
        paths.append(path)
    return paths


def _get_kakao_text_records(
    kakao_file_paths: list[str] | None,
) -> list[tuple[str, str]]:
    paths = _get_kakao_text_paths(kakao_file_paths)
    uncached_paths = [
        path for path in paths if str(Path(path)) not in _EXTRACTED_TEXT_CACHE
    ]
    outcomes = run_process_items_with_timeout(
        uncached_paths,
        _extract_text_in_process,
        timeout_seconds=DEFAULT_FILE_EXTRACTION_TIMEOUT_SECONDS,
        max_workers=FILE_EXTRACTION_MAX_WORKERS,
    )
    for path, (status, text) in zip(uncached_paths, outcomes):
        _EXTRACTED_TEXT_CACHE[str(Path(path))] = (
            str(text) if status == "ok" and text is not None else ""
        )
    return [
        (path, _EXTRACTED_TEXT_CACHE.get(str(Path(path)), ""))
        for path in paths
    ]


def _should_exclude_file(file: AnalyzedFile) -> bool:
    name = file.file_name.casefold()
    return name in SYSTEM_FILE_NAMES or file.is_hidden_or_system


def _select_latest_unique_files(
    files: list[tuple[AnalyzedFile, Path]],
    cancel_check: Callable[[], bool] | None = None,
) -> tuple[list[tuple[AnalyzedFile, Path]], list[dict]]:
    by_hash: dict[str, list[tuple[AnalyzedFile, Path]]] = {}
    excluded: list[dict] = []
    for item in files:
        # Full-content SHA-256 per file is the expensive part of this pass
        # (it reads every byte for dedup) — check for cancellation before
        # each one so a background worker can be interrupted promptly.
        _raise_if_cancelled(cancel_check)
        file_hash = _hash_file(item[1])
        if file_hash is None:
            excluded.append(_build_excluded_file_record(item[0], "read_error"))
            continue
        by_hash.setdefault(file_hash, []).append(item)

    selected: list[tuple[AnalyzedFile, Path]] = []
    for duplicates in by_hash.values():
        latest = max(duplicates, key=lambda item: item[0].modified_timestamp)
        selected.append(latest)
        for duplicate in duplicates:
            if duplicate is latest:
                continue
            excluded.append(_build_excluded_file_record(duplicate[0], "duplicate_content"))
    return selected, excluded


def _hash_file(path: Path) -> str | None:
    try:
        stat = path.stat()
    except OSError:
        return None
    cache_key = (str(path.resolve()), stat.st_size, stat.st_mtime_ns)
    cached = _FILE_HASH_CACHE.get(cache_key)
    if cached is not None:
        return cached

    digest = hashlib.sha256()
    try:
        with path.open("rb") as file:
            for block in iter(lambda: file.read(1024 * 1024), b""):
                digest.update(block)
    except OSError:
        return None
    value = digest.hexdigest()
    _FILE_HASH_CACHE[cache_key] = value
    return value


def _build_root_map(folder_paths: list[str]) -> dict[str, Path]:
    roots = [Path(path) for path in folder_paths]
    if len(roots) == 1:
        root = roots[0]
        return {root.name or str(root): root}

    root_map: dict[str, Path] = {}
    used: set[str] = set()
    for root in roots:
        base = root.name or str(root)
        namespace = base
        suffix = 2
        while namespace in used:
            namespace = f"{base} ({suffix})"
            suffix += 1
        used.add(namespace)
        root_map[namespace] = root
    return root_map


def _split_root_folder_paths(analysis_result: AnalysisResult) -> list[str]:
    return [
        path.strip()
        for path in analysis_result.root_folder_path.split("; ")
        if path.strip()
    ]


def _collect_rag_package_candidate_files(
    analysis_result: AnalysisResult,
    folder_paths: list[str],
    cancel_check: Callable[[], bool] | None = None,
) -> tuple[list[tuple[AnalyzedFile, Path]], list[dict]]:
    """Resolve, exclude, and de-duplicate files once using package rules."""
    root_map = _build_root_map(folder_paths)
    embeddable: list[tuple[AnalyzedFile, Path]] = []
    excluded_files: list[dict] = []
    for analyzed_file in analysis_result.all_files:
        _raise_if_cancelled(cancel_check)
        absolute_path = _resolve_absolute_path(analyzed_file, root_map)
        if _should_exclude_file(analyzed_file):
            excluded_files.append(
                _build_excluded_file_record(
                    analyzed_file,
                    "excluded_extension_or_system",
                )
            )
            continue
        if absolute_path is None or not absolute_path.is_file():
            excluded_files.append(
                _build_excluded_file_record(analyzed_file, "missing_file")
            )
            continue
        embeddable.append((analyzed_file, absolute_path))

    selected_files, duplicate_exclusions = _select_latest_unique_files(
        embeddable, cancel_check
    )
    excluded_files.extend(duplicate_exclusions)
    return selected_files, excluded_files


def _resolve_absolute_path(
    file: AnalyzedFile,
    root_map: dict[str, Path],
) -> Path | None:
    namespace, _, remainder = file.relative_path.partition("/")
    root = root_map.get(namespace)
    if root is None:
        return None
    return root / remainder if remainder else root / file.file_name


def _infer_source_type(file: AnalyzedFile) -> str:
    extension = Path(file.file_name).suffix.lower()
    if extension in {".eml", ".msg"}:
        return "email"
    if extension == ".txt" and "kakao" in file.file_name.casefold():
        return "kakao"
    return "document"


def _build_excluded_file_record(file: AnalyzedFile, reason: str) -> dict:
    return {
        "file_name": file.file_name,
        "relative_path": file.relative_path,
        "modified_at": file.modified_at,
        "reason": reason,
    }


def _build_analysis_id(analysis_result: AnalysisResult) -> str:
    raw = "|".join(
        [
            analysis_result.root_folder_path,
            str(analysis_result.total_file_count),
            str(analysis_result.total_size_bytes),
        ]
    )
    return hashlib.sha256(raw.encode("utf-8")).hexdigest()[:16]
